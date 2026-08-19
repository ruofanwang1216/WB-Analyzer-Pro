"""core/export_engine.py — PDF, TIFF, and PPTX figure export.

PDF export: uses QPdfWriter + QPainter.  All rendering is done here
independently of FigureCanvas (no scene reference needed).

PPTX export: uses python-pptx (OPTIONAL).  If the library is unavailable the
module still loads cleanly; only PPTX_AVAILABLE is set to False and callers
must check it before calling PPTXExporter.export().

TIFF export: writes an exact FigureCanvas snapshot as lossless LZW-compressed
RGB TIFF with publication-resolution DPI metadata.

Image cropping for PPTX is performed in Python (via QImage) BEFORE the image
is inserted into the presentation.  PowerPoint's crop tool is NOT used.
"""
from __future__ import annotations

import math
import os
import shutil
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path

import numpy as np
from PIL import Image
from PySide6.QtCore import QRect, QRectF, Qt
from PySide6.QtGui import (
    QColor, QFont, QImage, QPainter, QPen, QPixmap,
)

from core.image_transform import (
    apply_geometry_to_display,
    default_inverted_for_pil_image,
    geometry_transform_from_dict,
    image_array_to_uint16_luminance,
    image_transform_from_dict,
    transform_pixels_16_to_8,
)
from core.lane_composition import compose_lane_crops
from core.band_auto_fit import aspect_fit_placement
from core.layout_engine import (
    EXPORT_DPI, LayoutItem, LayoutResult, SCREEN_DPI, SCREEN_SCALE,
    emu_to_pt, pt_to_emu, pt_to_px,
)

# ── Optional python-pptx dependency ──────────────────────────────────────────
try:
    from pptx import Presentation as _Presentation
    from pptx.util import Emu as _Emu, Pt as _Pt
    from pptx.enum.text import (
        MSO_AUTO_SIZE as _MSO_AUTO_SIZE,
        MSO_VERTICAL_ANCHOR as _MSO_VERTICAL_ANCHOR,
        PP_ALIGN as _PP_ALIGN,
    )
    from pptx.dml.color import RGBColor as _RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ── Shared image utilities ────────────────────────────────────────────────────

# PowerPoint and Qt use different text rasterizers. These constants bias the
# editable PPTX output toward the live FigureCanvas preview.
_PPTX_FONT_SCALE = ((SCREEN_DPI / 72.0) / SCREEN_SCALE) * 0.78

# New presentations use the same 16:9 canvas as a standard widescreen
# PowerPoint slide. Figure content is never enlarged and is capped at roughly
# one-third of the slide, matching the supplied reference while leaving room
# for notes and other panels.
_PPTX_SLIDE_WIDTH_PT = 960.0
_PPTX_SLIDE_HEIGHT_PT = 540.0
_PPTX_LEFT_MARGIN_PT = 36.0
_PPTX_TOP_MARGIN_PT = 54.0
_PPTX_EDGE_MARGIN_PT = 36.0
_PPTX_FIGURE_MAX_WIDTH_RATIO = 0.36
_PPTX_FIGURE_MAX_HEIGHT_RATIO = 0.36


@dataclass(frozen=True)
class _PPTXPlacement:
    """Uniform PT-space transform that fits a figure on a slide."""

    scale: float
    offset_x_pt: float
    offset_y_pt: float

    def x(self, value_pt: float) -> float:
        return self.offset_x_pt + value_pt * self.scale

    def y(self, value_pt: float) -> float:
        return self.offset_y_pt + value_pt * self.scale


def _item_bounds_pt(item: LayoutItem) -> tuple[float, float, float, float]:
    """Return the visible axis-aligned bounds of one exported item in PT."""
    if item.kind == "line":
        x2 = item.x_pt + item.w_pt
        y2 = item.y_pt + item.h_pt
        inset = max(0.0, item.line_width_pt) / 2.0
        return (
            min(item.x_pt, x2) - inset,
            min(item.y_pt, y2) - inset,
            max(item.x_pt, x2) + inset,
            max(item.y_pt, y2) + inset,
        )

    left = min(item.x_pt, item.x_pt + item.w_pt)
    top = min(item.y_pt, item.y_pt + item.h_pt)
    width = abs(item.w_pt)
    height = abs(item.h_pt)
    if not item.rotation:
        return left, top, left + width, top + height

    radians = math.radians(item.rotation)
    rotated_w = abs(width * math.cos(radians)) + abs(height * math.sin(radians))
    rotated_h = abs(width * math.sin(radians)) + abs(height * math.cos(radians))
    center_x = left + width / 2.0
    center_y = top + height / 2.0
    return (
        center_x - rotated_w / 2.0,
        center_y - rotated_h / 2.0,
        center_x + rotated_w / 2.0,
        center_y + rotated_h / 2.0,
    )


def _fit_layout_to_slide(
    layout: LayoutResult,
    slide_width_pt: float,
    slide_height_pt: float,
) -> _PPTXPlacement:
    """Place content at natural size near the upper-left of a slide.

    Content is never enlarged. Oversized figures are uniformly reduced to the
    reference-sized box (roughly one-third of the slide) or the slide's safe
    area, whichever is smaller.
    """
    exported_items = [item for item in layout.items if item.kind != "divider"]
    if exported_items:
        bounds = [_item_bounds_pt(item) for item in exported_items]
        left = min(value[0] for value in bounds)
        top = min(value[1] for value in bounds)
        right = max(value[2] for value in bounds)
        bottom = max(value[3] for value in bounds)
    else:
        left = top = 0.0
        right = max(1.0, layout.canvas_width_pt)
        bottom = max(1.0, layout.canvas_height_pt)

    content_width = max(1.0, right - left)
    content_height = max(1.0, bottom - top)
    left_margin = min(_PPTX_LEFT_MARGIN_PT, slide_width_pt * 0.04)
    top_margin = min(_PPTX_TOP_MARGIN_PT, slide_height_pt * 0.10)
    right_margin = min(_PPTX_EDGE_MARGIN_PT, slide_width_pt * 0.04)
    bottom_margin = min(_PPTX_EDGE_MARGIN_PT, slide_height_pt * 0.07)
    available_width = min(
        max(1.0, slide_width_pt - left_margin - right_margin),
        max(1.0, slide_width_pt * _PPTX_FIGURE_MAX_WIDTH_RATIO),
    )
    available_height = min(
        max(1.0, slide_height_pt - top_margin - bottom_margin),
        max(1.0, slide_height_pt * _PPTX_FIGURE_MAX_HEIGHT_RATIO),
    )
    scale = min(
        1.0,
        available_width / content_width,
        available_height / content_height,
    )
    return _PPTXPlacement(
        scale=scale,
        offset_x_pt=left_margin - left * scale,
        offset_y_pt=top_margin - top * scale,
    )


def _fit_rect_to_slide(
    content_width_pt: float,
    content_height_pt: float,
    slide_width_pt: float,
    slide_height_pt: float,
) -> tuple[float, float, float, float]:
    """Return a natural-size upper-left rectangle for raster exports."""
    left = min(_PPTX_LEFT_MARGIN_PT, slide_width_pt * 0.04)
    top = min(_PPTX_TOP_MARGIN_PT, slide_height_pt * 0.10)
    right_margin = min(_PPTX_EDGE_MARGIN_PT, slide_width_pt * 0.04)
    bottom_margin = min(_PPTX_EDGE_MARGIN_PT, slide_height_pt * 0.07)
    available_width = min(
        max(1.0, slide_width_pt - left - right_margin),
        max(1.0, slide_width_pt * _PPTX_FIGURE_MAX_WIDTH_RATIO),
    )
    available_height = min(
        max(1.0, slide_height_pt - top - bottom_margin),
        max(1.0, slide_height_pt * _PPTX_FIGURE_MAX_HEIGHT_RATIO),
    )
    scale = min(
        1.0,
        available_width / max(1.0, content_width_pt),
        available_height / max(1.0, content_height_pt),
    )
    width = max(1.0, content_width_pt) * scale
    height = max(1.0, content_height_pt) * scale
    return (
        left,
        top,
        width,
        height,
    )


def _crop_array(pixels: np.ndarray, crop_px: dict | None) -> np.ndarray:
    if not crop_px:
        return pixels
    img_h, img_w = pixels.shape
    if img_w <= 0 or img_h <= 0:
        return pixels[0:0, 0:0]
    x = max(0, min(img_w - 1, int(round(float(crop_px.get("x", 0.0))))))
    y = max(0, min(img_h - 1, int(round(float(crop_px.get("y", 0.0))))))
    crop_w = max(1, int(round(float(crop_px.get("w", img_w - x)))))
    crop_h = max(1, int(round(float(crop_px.get("h", img_h - y)))))
    right = max(x + 1, min(img_w, x + crop_w))
    bottom = max(y + 1, min(img_h, y + crop_h))
    return pixels[y:bottom, x:right]


def _crop_qimage(
    image_path: str,
    crop_px: dict | None,
    image_transform: dict | None = None,
    lane_crops_px: list[dict] | None = None,
    geometry_transform: dict | None = None,
) -> QImage:
    """Load *image_path* and crop to *crop_px* (IMAGE_PX coordinates).

    Returns a null QImage if the file is missing or unreadable.
    The crop rect is applied in the source image's own pixel space.
    No coordinate-space conversion is performed here.
    """
    try:
        with Image.open(image_path) as img:
            default_inverted = default_inverted_for_pil_image(img, fallback=True)
            pixels = image_array_to_uint16_luminance(np.array(img))
    except Exception:
        return QImage()

    tone = image_transform_from_dict(
        image_transform,
        default_inverted=default_inverted,
    )
    geometry = geometry_transform_from_dict(geometry_transform)
    if geometry.is_identity():
        # The continuous crop is authoritative. Legacy Auto-Fit files may
        # still contain per-lane raw-space crops.
        if crop_px:
            pixels = _crop_array(pixels, crop_px)
        elif lane_crops_px:
            pixels = compose_lane_crops(pixels, lane_crops_px)
        if pixels.size == 0:
            return QImage()
        display = transform_pixels_16_to_8(pixels, tone)
    else:
        # Geometry belongs to presentation. Render it over the full immutable
        # source first, then apply the Canvas-space crop saved by the project.
        display = transform_pixels_16_to_8(pixels, tone)
        display = apply_geometry_to_display(display, geometry)
        display = _crop_array(display, crop_px)
        if display.size == 0:
            return QImage()
    display = np.ascontiguousarray(display)
    height, width = display.shape
    return QImage(
        display.data,
        width,
        height,
        display.strides[0],
        QImage.Format.Format_Grayscale8,
    ).copy()


def _qimage_to_png_bytes(image: QImage) -> bytes | None:
    if image.isNull():
        return None
    fd, tmp = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    try:
        image.save(tmp, "PNG")
        with open(tmp, "rb") as f:
            return f.read()
    finally:
        try:
            os.unlink(tmp)
        except OSError:
            pass


def _overwrite_file_contents(source_path: str, target_path: Path) -> None:
    """Copy file data into an existing inode and durably flush it.

    Keeping the inode preserves the target's permissions, ownership, Finder
    information, extended attributes, and resource forks. This also prevents
    Finder from briefly treating an updated PPTX as a deleted/recreated file.
    """
    with open(source_path, "rb") as source, open(target_path, "r+b") as target:
        target.seek(0)
        shutil.copyfileobj(source, target, length=1024 * 1024)
        target.truncate()
        target.flush()
        os.fsync(target.fileno())


def _save_presentation_safely(prs, output_path: str) -> None:  # type: ignore[no-untyped-def]
    """Save through a verified temp file while preserving an existing file.

    A new destination is installed atomically. For an existing destination,
    validated bytes are written into the original inode so filesystem metadata
    and Finder identity remain intact; a same-directory backup permits recovery
    if that final write or verification fails.
    """
    target = Path(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    fd, tmp = tempfile.mkstemp(
        prefix=f".{target.stem}-",
        suffix=".pptx",
        dir=str(target.parent),
    )
    os.close(fd)
    backup: str | None = None
    try:
        prs.save(tmp)
        if not zipfile.is_zipfile(tmp):
            raise RuntimeError("PPTX export produced an invalid PowerPoint package.")
        _Presentation(tmp)

        if not target.exists():
            os.replace(tmp, target)
            return

        backup_fd, backup = tempfile.mkstemp(
            prefix=f".{target.stem}-backup-",
            suffix=target.suffix,
            dir=str(target.parent),
        )
        os.close(backup_fd)
        shutil.copyfile(target, backup)
        try:
            _overwrite_file_contents(tmp, target)
            if not zipfile.is_zipfile(target):
                raise RuntimeError(
                    "PPTX export produced an invalid PowerPoint package."
                )
            _Presentation(target)
        except Exception:
            _overwrite_file_contents(backup, target)
            raise
    finally:
        for temporary_path in (tmp, backup):
            if not temporary_path:
                continue
            try:
                if Path(temporary_path).exists():
                    os.unlink(temporary_path)
            except OSError:
                pass


# ── PDF exporter ──────────────────────────────────────────────────────────────

class PDFExporter:
    """Renders a LayoutResult to a PDF file using QPdfWriter + QPainter.

    All coordinates are converted from PT to EXPORT_DPI pixels via pt_to_px().
    """

    def export(self, layout: LayoutResult, output_path: str) -> None:
        from PySide6.QtGui import QPdfWriter, QPageSize
        from PySide6.QtCore import QSizeF, QMarginsF

        writer = QPdfWriter(output_path)
        writer.setResolution(int(EXPORT_DPI))

        # Set page size to match the canvas (PT → mm: 1 pt = 25.4/72 mm)
        w_mm = layout.canvas_width_pt * 25.4 / 72.0
        h_mm = layout.canvas_height_pt * 25.4 / 72.0
        writer.setPageSize(QPageSize(QSizeF(w_mm, h_mm), QPageSize.Unit.Millimeter))
        writer.setPageMargins(QMarginsF(0.0, 0.0, 0.0, 0.0))

        painter = QPainter(writer)
        try:
            self._render(painter, layout, EXPORT_DPI)
        finally:
            painter.end()

    def export_image(
        self,
        image: QImage,
        canvas_width_pt: float,
        canvas_height_pt: float,
        output_path: str,
    ) -> None:
        """Write an exact raster snapshot of the figure to a one-page PDF."""
        if image.isNull():
            raise RuntimeError("PDF export could not render the figure snapshot.")

        from PySide6.QtGui import QPdfWriter, QPageSize
        from PySide6.QtCore import QSizeF, QMarginsF

        writer = QPdfWriter(output_path)
        writer.setResolution(int(EXPORT_DPI))

        w_mm = canvas_width_pt * 25.4 / 72.0
        h_mm = canvas_height_pt * 25.4 / 72.0
        writer.setPageSize(QPageSize(QSizeF(w_mm, h_mm), QPageSize.Unit.Millimeter))
        writer.setPageMargins(QMarginsF(0.0, 0.0, 0.0, 0.0))

        painter = QPainter(writer)
        try:
            target = writer.pageLayout().paintRectPixels(writer.resolution())
            if target.width() <= 0 or target.height() <= 0:
                target = QRect(0, 0, max(1, int(pt_to_px(canvas_width_pt, EXPORT_DPI))),
                               max(1, int(pt_to_px(canvas_height_pt, EXPORT_DPI))))
            painter.drawImage(target, image)
        finally:
            painter.end()

    # ── Rendering ─────────────────────────────────────────────────────────

    def _render(self, painter: QPainter, layout: LayoutResult, dpi: float) -> None:
        for item in sorted(layout.items, key=lambda i: i.z_order):
            self._draw_item(painter, item, dpi)

    def _draw_item(self, painter: QPainter, item: LayoutItem, dpi: float) -> None:
        x = pt_to_px(item.x_pt, dpi)
        y = pt_to_px(item.y_pt, dpi)
        w = pt_to_px(item.w_pt, dpi)
        h = pt_to_px(item.h_pt, dpi)
        rect = QRectF(x, y, max(w, 1.0), max(h, 1.0))
        irect = QRect(int(x), int(y), max(int(w), 1), max(int(h), 1))

        if item.kind == "blot":
            self._draw_blot(painter, item, irect, dpi)

        elif item.kind in ("label", "mw", "title", "panel_letter", "table_cell"):
            self._draw_text(painter, item, rect)

        elif item.kind == "line":
            pen = QPen(QColor(item.line_color or "#AAAAAA"))
            pen.setWidthF(pt_to_px(item.line_width_pt, dpi))
            painter.setPen(pen)
            painter.drawLine(int(x), int(y), int(x + w), int(y + h))

        elif item.kind == "divider":
            pen = QPen(QColor("#BBBBBB"))
            pen.setWidthF(pt_to_px(0.3, dpi))
            pen.setStyle(Qt.PenStyle.DashLine)
            painter.setPen(pen)
            painter.drawLine(int(x), int(y), int(x), int(y + h))

    def _draw_blot(
        self, painter: QPainter, item: LayoutItem, rect: QRect, dpi: float
    ) -> None:
        if item.image_path and Path(item.image_path).exists():
            img = _crop_qimage(
                item.image_path,
                item.image_crop_px,
                item.image_transform,
                item.image_lane_crops_px,
                geometry_transform=item.geometry_transform,
            )
            if not img.isNull():
                if item.preserve_image_aspect:
                    placement = aspect_fit_placement(
                        img.width(), img.height(), rect.width(), rect.height()
                    )
                    painter.drawImage(QRectF(
                        rect.x() + placement.x,
                        rect.y() + placement.y,
                        placement.width,
                        placement.height,
                    ), img)
                else:
                    painter.drawImage(QRectF(rect), img)
                return
        # Placeholder — grey rectangle with dashed border
        painter.fillRect(rect, QColor("#D8D8D8"))
        pen = QPen(QColor("#AAAAAA"))
        pen.setStyle(Qt.PenStyle.DashLine)
        pen.setWidthF(pt_to_px(0.5, dpi))
        painter.setPen(pen)
        painter.drawRect(rect)

    def _draw_text(
        self, painter: QPainter, item: LayoutItem, rect: QRectF
    ) -> None:
        font = QFont(item.font_family)
        font.setPointSizeF(item.font_size_pt)
        font.setBold(item.bold)
        font.setItalic(item.italic)
        font.setUnderline(item.underline)
        painter.setFont(font)
        painter.setPen(QColor("#000000"))

        h_flag = {
            "center": Qt.AlignmentFlag.AlignHCenter,
            "right":  Qt.AlignmentFlag.AlignRight,
        }.get(item.align, Qt.AlignmentFlag.AlignLeft)
        flags = h_flag | Qt.AlignmentFlag.AlignVCenter
        if item.rotation:
            painter.save()
            try:
                center = rect.center()
                painter.translate(center)
                painter.rotate(item.rotation)
                painter.translate(-center)
                painter.drawText(rect, flags, item.text)
                painter.setOpacity(0.28)
                painter.translate(0.18, 0.0)
                painter.drawText(rect, flags, item.text)
            finally:
                painter.restore()
            return
        painter.drawText(rect, flags, item.text)


# ── TIFF exporter ─────────────────────────────────────────────────────────────

class TIFFExporter:
    """Writes a lossless, publication-resolution raster snapshot as TIFF."""

    DEFAULT_DPI = int(EXPORT_DPI)

    @staticmethod
    def render_scale_for_dpi(dpi: float = EXPORT_DPI) -> float:
        """Return the FigureCanvas scene-render scale for the requested DPI."""
        dpi = float(dpi)
        if dpi <= 0.0:
            raise ValueError("TIFF DPI must be greater than zero.")
        return dpi / (72.0 * SCREEN_SCALE)

    def export_image(
        self,
        image: QImage,
        output_path: str,
        *,
        dpi: float = EXPORT_DPI,
    ) -> None:
        """Save *image* as an RGB TIFF with lossless LZW compression."""
        if image.isNull():
            raise RuntimeError("TIFF export could not render the figure snapshot.")
        dpi = float(dpi)
        if dpi <= 0.0:
            raise ValueError("TIFF DPI must be greater than zero.")

        png_bytes = _qimage_to_png_bytes(image)
        if not png_bytes:
            raise RuntimeError("TIFF export could not encode the figure snapshot.")

        import io

        target = Path(output_path)
        target.parent.mkdir(parents=True, exist_ok=True)
        fd, tmp = tempfile.mkstemp(
            prefix=f".{target.stem}-",
            suffix=".tiff",
            dir=str(target.parent),
        )
        os.close(fd)
        try:
            with Image.open(io.BytesIO(png_bytes)) as source:
                source.convert("RGB").save(
                    tmp,
                    format="TIFF",
                    compression="tiff_lzw",
                    dpi=(dpi, dpi),
                )
            with Image.open(tmp) as check:
                check.load()
                if check.format != "TIFF" or check.size != (
                    image.width(),
                    image.height(),
                ):
                    raise RuntimeError("TIFF export produced an invalid image.")
            os.replace(tmp, target)
        finally:
            try:
                if Path(tmp).exists():
                    os.unlink(tmp)
            except OSError:
                pass


# ── PPTX exporter ─────────────────────────────────────────────────────────────

class PPTXExporter:
    """Renders a LayoutResult to a PPTX file.

    Each blot → add_picture() with Python-cropped image (no PowerPoint crop).
    Each label → add_textbox() (independently editable in PowerPoint).
    Each line  → add_connector().
    """

    def export(self, layout: LayoutResult, output_path: str) -> None:
        if not PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx is not installed.  Run: pip install python-pptx"
            )

        prs = _Presentation()
        slide_w = pt_to_emu(_PPTX_SLIDE_WIDTH_PT)
        slide_h = pt_to_emu(_PPTX_SLIDE_HEIGHT_PT)
        prs.slide_width = slide_w
        prs.slide_height = slide_h

        layout_obj = prs.slide_layouts[6]   # blank layout
        slide = prs.slides.add_slide(layout_obj)

        self._add_layout_items(
            slide,
            layout,
            emu_to_pt(prs.slide_width),
            emu_to_pt(prs.slide_height),
        )

        _save_presentation_safely(prs, output_path)

    def export_append_slide(self, layout: LayoutResult, existing_path: str) -> None:
        """Append a new slide to an existing PPTX file and save in place."""
        if not PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx is not installed.  Run: pip install python-pptx"
            )

        prs = _Presentation(existing_path)
        layout_obj = prs.slide_layouts[6]   # blank layout
        slide = prs.slides.add_slide(layout_obj)

        self._add_layout_items(
            slide,
            layout,
            emu_to_pt(prs.slide_width),
            emu_to_pt(prs.slide_height),
        )

        _save_presentation_safely(prs, existing_path)

    def export_image(
        self,
        image: QImage,
        canvas_width_pt: float,
        canvas_height_pt: float,
        output_path: str,
    ) -> None:
        """Export an exact raster snapshot of the figure as a PPTX slide."""
        if not PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx is not installed.  Run: pip install python-pptx"
            )
        png_bytes = _qimage_to_png_bytes(image)
        if not png_bytes:
            raise RuntimeError("PPTX export could not render the figure snapshot.")

        prs = _Presentation()
        prs.slide_width = pt_to_emu(_PPTX_SLIDE_WIDTH_PT)
        prs.slide_height = pt_to_emu(_PPTX_SLIDE_HEIGHT_PT)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_snapshot_picture(
            slide,
            png_bytes,
            canvas_width_pt,
            canvas_height_pt,
            emu_to_pt(prs.slide_width),
            emu_to_pt(prs.slide_height),
        )
        _save_presentation_safely(prs, output_path)

    def export_append_image(
        self,
        image: QImage,
        canvas_width_pt: float,
        canvas_height_pt: float,
        existing_path: str,
    ) -> None:
        """Append an exact raster snapshot of the figure to an existing PPTX."""
        if not PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx is not installed.  Run: pip install python-pptx"
            )
        png_bytes = _qimage_to_png_bytes(image)
        if not png_bytes:
            raise RuntimeError("PPTX export could not render the figure snapshot.")

        prs = _Presentation(existing_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_snapshot_picture(
            slide,
            png_bytes,
            canvas_width_pt,
            canvas_height_pt,
            emu_to_pt(prs.slide_width),
            emu_to_pt(prs.slide_height),
        )
        _save_presentation_safely(prs, existing_path)

    @staticmethod
    def _add_snapshot_picture(
        slide,
        png_bytes: bytes,
        canvas_width_pt: float,
        canvas_height_pt: float,
        slide_width_pt: float,
        slide_height_pt: float,
    ):
        import io
        left, top, width, height = _fit_rect_to_slide(
            canvas_width_pt,
            canvas_height_pt,
            slide_width_pt,
            slide_height_pt,
        )
        return slide.shapes.add_picture(
            io.BytesIO(png_bytes),
            _Emu(pt_to_emu(left)),
            _Emu(pt_to_emu(top)),
            _Emu(pt_to_emu(width)),
            _Emu(pt_to_emu(height)),
        )

    # ── Item dispatch ──────────────────────────────────────────────────────

    def _add_layout_items(
        self,
        slide,
        layout: LayoutResult,
        slide_width_pt: float,
        slide_height_pt: float,
    ) -> None:  # type: ignore[type-arg]
        placement = _fit_layout_to_slide(layout, slide_width_pt, slide_height_pt)
        for item in sorted(layout.items, key=lambda i: i.z_order):
            self._add_item(slide, item, placement)

    def _add_item(  # type: ignore[type-arg,no-untyped-def]
        self,
        slide,
        item: LayoutItem,
        placement: _PPTXPlacement,
    ):
        left = _Emu(pt_to_emu(placement.x(item.x_pt)))
        top = _Emu(pt_to_emu(placement.y(item.y_pt)))
        width = _Emu(pt_to_emu(item.w_pt * placement.scale))
        height = _Emu(pt_to_emu(item.h_pt * placement.scale))

        if item.kind == "blot":
            return self._add_blot(
                slide,
                item,
                left,
                top,
                width,
                height,
                placement.scale,
            )

        elif item.kind in ("label", "mw", "title", "panel_letter", "table_cell"):
            return self._add_textbox(
                slide,
                item,
                left,
                top,
                width,
                height,
                placement.scale,
            )

        elif item.kind == "line":
            return self._add_line(slide, item, placement)

        # dividers are omitted from PPTX (visual noise at publication scale)
        return None

    def _add_blot(
        self,
        slide,
        item: LayoutItem,
        left,
        top,
        width,
        height,
        scale: float,
    ) -> None:
        if item.image_path and Path(item.image_path).exists():
            image = _crop_qimage(
                item.image_path,
                item.image_crop_px,
                item.image_transform,
                item.image_lane_crops_px,
                geometry_transform=item.geometry_transform,
            )
            png_bytes = _qimage_to_png_bytes(image)
            if png_bytes:
                import io
                buf = io.BytesIO(png_bytes)
                if item.preserve_image_aspect:
                    placement = aspect_fit_placement(
                        image.width(), image.height(), width, height
                    )
                    picture_left = left + int(round(placement.x))
                    picture_top = top + int(round(placement.y))
                    picture_width = int(round(placement.width))
                    picture_height = int(round(placement.height))
                else:
                    picture_left = left
                    picture_top = top
                    picture_width = width
                    picture_height = height
                picture = slide.shapes.add_picture(
                    buf,
                    picture_left,
                    picture_top,
                    picture_width,
                    picture_height,
                )
                self._apply_blot_border(picture, scale)
                return picture
        # Placeholder rectangle when no image is available
        from pptx.dml.color import RGBColor as _RGB2
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _RGB2(0xD0, 0xD0, 0xD0)
        self._apply_blot_border(shape, scale)
        return shape

    @staticmethod
    def _apply_blot_border(shape, scale: float = 1.0) -> None:  # type: ignore[no-untyped-def]
        shape.line.color.rgb = _RGBColor(0x00, 0x00, 0x00)
        shape.line.width = _Pt(max(0.25, scale))

    def _add_textbox(
        self,
        slide,
        item: LayoutItem,
        left,
        top,
        width,
        height,
        scale: float,
    ) -> None:
        txb = slide.shapes.add_textbox(left, top, width, height)
        txb.rotation = item.rotation
        tf = txb.text_frame
        tf.word_wrap = False
        tf.auto_size = _MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = _MSO_VERTICAL_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.space_before = _Pt(0)
        p.space_after = _Pt(0)
        p.line_spacing = 1.0

        # Alignment
        align_map = {
            "center": _PP_ALIGN.CENTER,
            "right":  _PP_ALIGN.RIGHT,
            "left":   _PP_ALIGN.LEFT,
        }
        p.alignment = align_map.get(item.align, _PP_ALIGN.LEFT)

        run = p.add_run()
        run.text = item.text
        run.font.name = item.font_family
        run.font.size = _Pt(item.font_size_pt * _PPTX_FONT_SCALE * scale)
        run.font.bold = item.bold
        run.font.italic = item.italic
        run.font.underline = item.underline
        run.font.color.rgb = _RGBColor(0x00, 0x00, 0x00)
        return txb

    def _add_line(
        self,
        slide,
        item: LayoutItem,
        placement: _PPTXPlacement,
    ) -> None:
        from pptx.util import Emu as _E2
        x1 = _E2(pt_to_emu(placement.x(item.x_pt)))
        y1 = _E2(pt_to_emu(placement.y(item.y_pt)))
        x2 = _E2(pt_to_emu(placement.x(item.x_pt + item.w_pt)))
        y2 = _E2(pt_to_emu(placement.y(item.y_pt + item.h_pt)))
        connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
        color = QColor(item.line_color or "#AAAAAA")
        connector.line.color.rgb = _RGBColor(color.red(), color.green(), color.blue())
        connector.line.width = _Emu(
            pt_to_emu(item.line_width_pt * placement.scale)
        )
        # PowerPoint otherwise inherits the theme's effect reference, which can
        # render as a grey shadow below horizontal connector lines.
        connector.shadow.inherit = False
        return connector
