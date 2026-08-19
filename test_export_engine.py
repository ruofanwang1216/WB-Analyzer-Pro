import unittest
import os
import stat
import sys
import zipfile
from tempfile import TemporaryDirectory

import numpy as np
from PIL import Image
from PySide6.QtGui import QColor, QImage

from core.export_engine import (
    PDFExporter,
    PPTX_AVAILABLE,
    PPTXExporter,
    TIFFExporter,
    _crop_qimage,
    _fit_layout_to_slide,
    _fit_rect_to_slide,
    _PPTX_SLIDE_HEIGHT_PT,
    _PPTX_SLIDE_WIDTH_PT,
)
from core.figure_project import SourceRef
from core.layout_engine import EMU_PER_PT, LayoutItem, LayoutResult, SCREEN_DPI, SCREEN_SCALE


if PPTX_AVAILABLE:
    from pptx import Presentation


class LaneCropExportTests(unittest.TestCase):
    def test_geometry_is_rendered_before_canvas_space_crop(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/source.png"
            pixels = np.array([[10, 20, 200, 240], [11, 21, 201, 241]], dtype=np.uint8)
            Image.fromarray(pixels).save(path)
            image = _crop_qimage(
                path,
                {"x": 0, "y": 0, "w": 2, "h": 2},
                {"low": 0, "high": 65535, "gamma": 1.0, "inverted": False},
                geometry_transform={"rotation": 0.0, "flip_x": True, "flip_y": False},
            )

        self.assertEqual(image.width(), 2)
        self.assertEqual(image.height(), 2)
        self.assertEqual(
            [[image.pixelColor(x, y).red() for x in range(2)] for y in range(2)],
            [[240, 200], [241, 201]],
        )

    def test_export_composes_equal_lane_crops_before_scaling(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/source.png"
            Image.fromarray(np.full((30, 50), 180, dtype=np.uint8)).save(path)
            image = _crop_qimage(
                path,
                None,
                {"low": 0, "high": 65535, "gamma": 1.0, "inverted": False},
                [
                    {"x": 2, "y": 4, "w": 12, "h": 9},
                    {"x": 28, "y": 14, "w": 12, "h": 9},
                ],
            )
        self.assertEqual(image.width(), 24)
        self.assertEqual(image.height(), 9)

    def test_continuous_crop_takes_priority_over_legacy_lane_crops(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/source.png"
            pixels = np.zeros((20, 30), dtype=np.uint16)
            pixels[5:7, 2:5] = 65535
            pixels[8:10, 12:15] = 65535
            pixels[11:13, 22:25] = 65535
            Image.fromarray(pixels).save(path)
            image = _crop_qimage(
                path,
                {"x": 0, "y": 0, "w": 30, "h": 20},
                {"low": 0, "high": 65535, "gamma": 1.0, "inverted": False},
                [
                    {"x": 0, "y": 5, "w": 10, "h": 2},
                    {"x": 10, "y": 8, "w": 10, "h": 2},
                    {"x": 20, "y": 11, "w": 10, "h": 2},
                ],
            )
        self.assertEqual(image.width(), 30)
        self.assertEqual(image.height(), 20)
        self.assertEqual(
            [image.pixelColor(x, y).red() for x, y in ((3, 5), (13, 8), (23, 11))],
            [255, 255, 255],
        )
        self.assertEqual(
            [image.pixelColor(x, 5).red() for x in (13, 23)],
            [0, 0],
        )


@unittest.skipUnless(PPTX_AVAILABLE, "python-pptx is not installed")
class PPTXExporterTests(unittest.TestCase):
    _expected_font_scale = ((SCREEN_DPI / 72.0) / SCREEN_SCALE) * 0.78

    def _layout(self, text: str = "HEK293T") -> LayoutResult:
        return LayoutResult(
            canvas_width_pt=456.0,
            canvas_height_pt=220.0,
            items=[
                LayoutItem(
                    kind="title",
                    x_pt=160.0,
                    y_pt=20.0,
                    w_pt=42.0,
                    h_pt=18.0,
                    text=text,
                    font_size_pt=16.0,
                    bold=True,
                    align="center",
                    z_order=2,
                ),
                LayoutItem(
                    kind="table_cell",
                    x_pt=40.0,
                    y_pt=50.0,
                    w_pt=48.0,
                    h_pt=17.0,
                    text="pcDNA-SMYD3",
                    font_size_pt=13.0,
                    bold=True,
                    align="right",
                    z_order=2,
                    source_ref=SourceRef(panel_idx=0, table_row=1, table_col=0, field="condition_cell"),
                ),
                LayoutItem(
                    kind="table_cell",
                    x_pt=210.0,
                    y_pt=56.0,
                    w_pt=18.0,
                    h_pt=30.0,
                    text="1000",
                    font_size_pt=13.0,
                    align="center",
                    z_order=2,
                    source_ref=SourceRef(panel_idx=0, table_row=1, table_col=1, field="condition_cell"),
                ),
                LayoutItem(
                    kind="blot",
                    x_pt=96.0,
                    y_pt=80.0,
                    w_pt=240.0,
                    h_pt=18.0,
                    z_order=1,
                ),
                LayoutItem(
                    kind="line",
                    x_pt=110.0,
                    y_pt=42.0,
                    w_pt=180.0,
                    h_pt=0.0,
                    line_color="#222222",
                    line_width_pt=1.0,
                    z_order=3,
                ),
            ],
        )

    def test_export_text_boxes_do_not_wrap(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/figure.pptx"

            PPTXExporter().export(self._layout(), path)

            prs = Presentation(path)
            text_shapes = [
                shape for shape in prs.slides[0].shapes
                if getattr(shape, "has_text_frame", False) and shape.text
            ]
            table_shapes = [
                shape for shape in prs.slides[0].shapes
                if getattr(shape, "has_table", False)
            ]
            self.assertEqual(len(text_shapes), 3)
            self.assertEqual(len(table_shapes), 0)
            self.assertTrue(all(shape.text_frame.word_wrap is False for shape in text_shapes))
            self.assertEqual([shape.text for shape in text_shapes], ["HEK293T", "pcDNA-SMYD3", "1000"])
            self.assertAlmostEqual(
                text_shapes[0].text_frame.paragraphs[0].runs[0].font.size.pt,
                16.0
                * self._expected_font_scale
                * _fit_layout_to_slide(
                    self._layout(),
                    _PPTX_SLIDE_WIDTH_PT,
                    _PPTX_SLIDE_HEIGHT_PT,
                ).scale,
                places=1,
            )
            placement = _fit_layout_to_slide(
                self._layout(),
                _PPTX_SLIDE_WIDTH_PT,
                _PPTX_SLIDE_HEIGHT_PT,
            )
            self.assertEqual(
                text_shapes[0].top,
                int(placement.y(20.0) * EMU_PER_PT),
            )
            self.assertEqual(
                text_shapes[1].top,
                int(placement.y(50.0) * EMU_PER_PT),
            )
            self.assertEqual(
                text_shapes[2].top,
                int(placement.y(56.0) * EMU_PER_PT),
            )
            line_shape = next(
                shape for shape in prs.slides[0].shapes
                if hasattr(shape, "line")
                and str(getattr(shape.line.color, "rgb", "")) == "222222"
            )
            self.assertEqual(str(line_shape.line.color.rgb), "222222")
            self.assertFalse(line_shape.shadow.inherit)
            blot_shape = prs.slides[0].shapes[0]
            self.assertEqual(str(blot_shape.line.color.rgb), "000000")

    def test_condition_cells_export_at_layout_item_positions(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/figure.pptx"

            PPTXExporter().export(self._layout(), path)

            prs = Presentation(path)
            self.assertEqual(
                sum(1 for shape in prs.slides[0].shapes if getattr(shape, "has_table", False)),
                0,
            )
            table_text_shapes = [
                shape for shape in prs.slides[0].shapes
                if getattr(shape, "has_text_frame", False)
                and shape.text in {"pcDNA-SMYD3", "1000"}
            ]
            self.assertEqual(len(table_text_shapes), 2)
            placement = _fit_layout_to_slide(
                self._layout(),
                _PPTX_SLIDE_WIDTH_PT,
                _PPTX_SLIDE_HEIGHT_PT,
            )
            self.assertEqual(
                table_text_shapes[0].left,
                int(placement.x(40.0) * EMU_PER_PT),
            )
            self.assertEqual(
                table_text_shapes[1].left,
                int(placement.x(210.0) * EMU_PER_PT),
            )
            self.assertEqual(
                table_text_shapes[0].top,
                int(placement.y(50.0) * EMU_PER_PT),
            )
            self.assertEqual(
                table_text_shapes[1].top,
                int(placement.y(56.0) * EMU_PER_PT),
            )

    def test_new_slide_uses_widescreen_and_keeps_natural_size_near_upper_left(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/centred.pptx"
            layout = self._layout()
            layout.items.append(
                LayoutItem(
                    kind="title",
                    x_pt=120.0,
                    y_pt=-24.0,
                    w_pt=110.0,
                    h_pt=18.0,
                    text="48h",
                    font_size_pt=15.0,
                    align="center",
                    z_order=4,
                )
            )

            PPTXExporter().export(layout, path)

            prs = Presentation(path)
            slide = prs.slides[0]
            shapes = list(slide.shapes)
            left = min(shape.left for shape in shapes)
            top = min(shape.top for shape in shapes)
            right = max(shape.left + shape.width for shape in shapes)
            bottom = max(shape.top + shape.height for shape in shapes)
            self.assertGreaterEqual(left, 0)
            self.assertGreaterEqual(top, 0)
            self.assertLessEqual(right, prs.slide_width)
            self.assertLessEqual(bottom, prs.slide_height)
            self.assertEqual(prs.slide_width, int(_PPTX_SLIDE_WIDTH_PT * EMU_PER_PT))
            self.assertEqual(prs.slide_height, int(_PPTX_SLIDE_HEIGHT_PT * EMU_PER_PT))
            placement = _fit_layout_to_slide(layout, 960.0, 540.0)
            self.assertEqual(placement.scale, 1.0)
            self.assertAlmostEqual(left, 36.0 * EMU_PER_PT, delta=2.0 * EMU_PER_PT)
            self.assertAlmostEqual(top, 54.0 * EMU_PER_PT, delta=2.0 * EMU_PER_PT)

    def test_append_slide_preserves_existing_file_identity_and_metadata(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/existing.pptx"
            exporter = PPTXExporter()
            exporter.export(self._layout("Original"), path)
            os.chmod(path, 0o640)
            inode_before = os.stat(path).st_ino
            mode_before = stat.S_IMODE(os.stat(path).st_mode)
            xattr_name = (
                b"com.wb-analyzer.export-test"
                if sys.platform == "darwin"
                else b"user.wb-analyzer.export-test"
            )
            xattr_supported = hasattr(os, "setxattr")
            if xattr_supported:
                try:
                    os.setxattr(path, xattr_name, b"keep-me")
                except OSError:
                    xattr_supported = False

            exporter.export_append_slide(self._layout("Appended"), path)

            self.assertTrue(zipfile.is_zipfile(path))
            self.assertEqual(os.stat(path).st_ino, inode_before)
            self.assertEqual(stat.S_IMODE(os.stat(path).st_mode), mode_before)
            if xattr_supported:
                self.assertEqual(os.getxattr(path, xattr_name), b"keep-me")
            prs = Presentation(path)
            self.assertEqual(len(prs.slides), 2)
            slide_text = "\n".join(shape.text for shape in prs.slides[-1].shapes if getattr(shape, "has_text_frame", False))
            self.assertIn("Appended", slide_text)

    def test_append_slide_preserves_existing_format_parts_byte_for_byte(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/formatted.pptx"
            prs = Presentation()
            prs.slide_width = int(720.0 * EMU_PER_PT)
            prs.slide_height = int(405.0 * EMU_PER_PT)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Existing formatted slide"
            prs.save(path)

            required_package_updates = {
                "[Content_Types].xml",
                "ppt/presentation.xml",
                "ppt/_rels/presentation.xml.rels",
            }
            with zipfile.ZipFile(path) as archive:
                before = {
                    name: archive.read(name)
                    for name in archive.namelist()
                    if name not in required_package_updates
                }

            PPTXExporter().export_append_slide(self._layout(), path)

            with zipfile.ZipFile(path) as archive:
                after = {
                    name: archive.read(name)
                    for name in before
                }
            self.assertEqual(after, before)
            result = Presentation(path)
            self.assertEqual(result.slide_width, int(720.0 * EMU_PER_PT))
            self.assertEqual(result.slide_height, int(405.0 * EMU_PER_PT))

    def test_append_slide_keeps_natural_size_near_upper_left(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/existing-wide.pptx"
            prs = Presentation()
            prs.slide_width = int(720.0 * EMU_PER_PT)
            prs.slide_height = int(405.0 * EMU_PER_PT)
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(path)
            layout = self._layout()
            layout.items[0].y_pt = -30.0

            PPTXExporter().export_append_slide(layout, path)

            result = Presentation(path)
            shapes = list(result.slides[-1].shapes)
            left = min(shape.left for shape in shapes)
            top = min(shape.top for shape in shapes)
            right = max(shape.left + shape.width for shape in shapes)
            bottom = max(shape.top + shape.height for shape in shapes)
            self.assertGreaterEqual(left, 0)
            self.assertGreaterEqual(top, 0)
            self.assertLessEqual(right, result.slide_width)
            self.assertLessEqual(bottom, result.slide_height)
            placement = _fit_layout_to_slide(layout, 720.0, 405.0)
            self.assertLess(placement.scale, 1.0)
            self.assertAlmostEqual(left, 28.8 * EMU_PER_PT, delta=2.0 * EMU_PER_PT)
            self.assertAlmostEqual(top, 40.5 * EMU_PER_PT, delta=2.0 * EMU_PER_PT)
            self.assertLessEqual(
                right - left,
                720.0 * 0.36 * EMU_PER_PT + 2.0 * EMU_PER_PT,
            )

    def test_export_image_writes_single_exact_snapshot_picture(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/snapshot.pptx"
            image = QImage(200, 100, QImage.Format.Format_ARGB32)
            image.fill(QColor("#FFFFFF"))

            PPTXExporter().export_image(image, 456.0, 220.0, path)

            prs = Presentation(path)
            self.assertEqual(prs.slide_width, int(round(_PPTX_SLIDE_WIDTH_PT * EMU_PER_PT)))
            self.assertEqual(prs.slide_height, int(round(_PPTX_SLIDE_HEIGHT_PT * EMU_PER_PT)))
            shapes = list(prs.slides[0].shapes)
            self.assertEqual(len(shapes), 1)
            left, top, width, height = _fit_rect_to_slide(
                456.0,
                220.0,
                _PPTX_SLIDE_WIDTH_PT,
                _PPTX_SLIDE_HEIGHT_PT,
            )
            self.assertEqual(shapes[0].left, int(left * EMU_PER_PT))
            self.assertEqual(shapes[0].top, int(top * EMU_PER_PT))
            self.assertEqual(shapes[0].width, int(width * EMU_PER_PT))
            self.assertEqual(shapes[0].height, int(height * EMU_PER_PT))
            self.assertAlmostEqual(
                shapes[0].width / prs.slide_width,
                0.36,
                places=3,
            )


class PDFExporterTests(unittest.TestCase):
    def test_export_image_writes_pdf_snapshot(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/snapshot.pdf"
            image = QImage(200, 100, QImage.Format.Format_ARGB32)
            image.fill(QColor("#FFFFFF"))

            PDFExporter().export_image(image, 456.0, 220.0, path)

            with open(path, "rb") as handle:
                self.assertEqual(handle.read(5), b"%PDF-")


class TIFFExporterTests(unittest.TestCase):
    def test_export_image_writes_lossless_300_dpi_tiff(self) -> None:
        with TemporaryDirectory() as tmp:
            path = f"{tmp}/figure.tiff"
            image = QImage(200, 100, QImage.Format.Format_ARGB32)
            image.fill(QColor("#F4F4F4"))

            TIFFExporter().export_image(image, path)

            with Image.open(path) as exported:
                exported.load()
                self.assertEqual(exported.format, "TIFF")
                self.assertEqual(exported.mode, "RGB")
                self.assertEqual(exported.size, (200, 100))
                self.assertEqual(exported.tag_v2.get(259), 5)  # LZW
                self.assertAlmostEqual(float(exported.tag_v2.get(282)), 300.0)
                self.assertAlmostEqual(float(exported.tag_v2.get(283)), 300.0)
                self.assertEqual(exported.tag_v2.get(296), 2)  # inches

    def test_render_scale_matches_figure_canvas_scene_units(self) -> None:
        self.assertAlmostEqual(
            TIFFExporter.render_scale_for_dpi(300.0),
            300.0 / (72.0 * SCREEN_SCALE),
        )

    def test_export_rejects_null_image(self) -> None:
        with TemporaryDirectory() as tmp:
            with self.assertRaises(RuntimeError):
                TIFFExporter().export_image(QImage(), f"{tmp}/empty.tiff")


if __name__ == "__main__":
    unittest.main()
